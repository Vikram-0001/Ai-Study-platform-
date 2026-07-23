import os
from typing import List, Dict, Any
from pypdf import PdfReader
from pptx import Presentation
from docx import Document as DocxDocument

class DocumentParser:
    def parse_file(self, file_path: str, file_type: str) -> List[Dict[str, Any]]:
        """
        Parses different file formats into content pages/segments.
        Each page contains raw text content and the respective page index.
        Handles both local paths and remote URLs (downloads remote files to temporary storage).
        """
        import tempfile
        import httpx
        
        is_url = file_path.startswith("http://") or file_path.startswith("https://")
        temp_file_path = file_path
        
        if is_url:
            try:
                suffix = f".{file_type.lower()}"
                with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                    temp_file_path = temp_file.name
                    with httpx.Client() as client:
                        response = client.get(file_path)
                        if response.status_code == 200:
                            temp_file.write(response.content)
                        else:
                            raise ValueError(f"Failed to download remote file from {file_path}, status: {response.status_code}")
            except Exception as e:
                print(f"Error downloading remote file {file_path}: {e}")
                # clean up temp file if created
                if temp_file_path != file_path and os.path.exists(temp_file_path):
                    import os
                    os.remove(temp_file_path)
                return []
                
        try:
            file_type = file_type.lower()
            if file_type == "pdf":
                pages = self._parse_pdf(temp_file_path)
            elif file_type == "pptx":
                pages = self._parse_pptx(temp_file_path)
            elif file_type in ["docx", "doc"]:
                pages = self._parse_docx(temp_file_path)
            elif file_type in ["txt", "md"]:
                pages = self._parse_txt(temp_file_path)
            else:
                try:
                    pages = self._parse_txt(temp_file_path)
                except Exception:
                    raise ValueError(f"Unsupported file format extension: {file_type}")
        finally:
            if is_url and temp_file_path != file_path and os.path.exists(temp_file_path):
                import os
                os.remove(temp_file_path)
                
        return pages

    def _parse_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        pages = []
        try:
            reader = PdfReader(file_path)
            for idx, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                pages.append({
                    "text": text.strip(),
                    "page_number": idx + 1
                })
        except Exception as e:
            print(f"Error parsing PDF file {file_path}: {e}")
        return pages

    def _parse_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        slides = []
        try:
            prs = Presentation(file_path)
            for idx, slide in enumerate(prs.slides):
                text_runs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_runs.append(shape.text)
                text = "\n".join(text_runs)
                slides.append({
                    "text": text.strip(),
                    "page_number": idx + 1
                })
        except Exception as e:
            print(f"Error parsing PPTX slides {file_path}: {e}")
        return slides

    def _parse_docx(self, file_path: str) -> List[Dict[str, Any]]:
        pages = []
        try:
            doc = DocxDocument(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            
            # Since Word docs don't have distinct "pages" naturally without page break indicators,
            # we group paragraphs into arbitrary virtual pages (e.g. 5 paragraphs per page)
            chunk_paragraphs = 5
            for idx in range(0, len(paragraphs), chunk_paragraphs):
                text = "\n".join(paragraphs[idx : idx + chunk_paragraphs])
                pages.append({
                    "text": text.strip(),
                    "page_number": (idx // chunk_paragraphs) + 1
                })
        except Exception as e:
            print(f"Error parsing DOCX file {file_path}: {e}")
        return pages

    def _parse_txt(self, file_path: str) -> List[Dict[str, Any]]:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            # Virtual split by character length to create "pages" (approx 2000 chars per page)
            page_size = 2000
            pages = []
            for idx, start in enumerate(range(0, len(content), page_size)):
                text = content[start : start + page_size]
                pages.append({
                    "text": text.strip(),
                    "page_number": idx + 1
                })
            return pages
        except Exception as e:
            print(f"Error parsing TXT file {file_path}: {e}")
            return []

    def chunk_document(self, pages: List[Dict[str, Any]], chunk_size: int = 1000, chunk_overlap: int = 200) -> List[Dict[str, Any]]:
        """
        Chunks list of pages into smaller units using a sliding text window.
        Preserves reference page numbers and content flow.
        """
        chunks = []
        chunk_index = 0

        for page in pages:
            text = page["text"]
            page_num = page["page_number"]
            
            if not text:
                continue
                
            # If text fits, keep it as single chunk
            if len(text) <= chunk_size:
                chunks.append({
                    "text": text,
                    "page_number": page_num,
                    "chunk_id": chunk_index
                })
                chunk_index += 1
                continue
            
            # Slide a window across page text
            start = 0
            while start < len(text):
                end = start + chunk_size
                chunk_text = text[start:end]
                chunks.append({
                    "text": chunk_text.strip(),
                    "page_number": page_num,
                    "chunk_id": chunk_index
                })
                chunk_index += 1
                start += (chunk_size - chunk_overlap)
                
        return chunks

document_parser = DocumentParser()
